"""Generate and exercise a fully offline mixed-format Atlas batch."""

import argparse
import hashlib
import json
import os
import sys
import tempfile
import time
from contextlib import chdir
from pathlib import Path
from unittest.mock import patch


ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "src"
if str(SRC) not in sys.path:
    sys.path.insert(0, str(SRC))

# This validation tool must be safe even when the caller forgets the shell
# guards. It never permits dotenv loading or a real model credential.
os.environ["PYTHON_DOTENV_DISABLED"] = "1"
os.environ["OPENAI_API_KEY"] = "bulk-ingestion-offline-placeholder"

import products.atlas.concept_enrichment as concept_enrichment  # noqa: E402
import products.atlas.intake_service as intake_service  # noqa: E402
import wingman.core.concept_registry_storage as concept_registry_storage  # noqa: E402
import wingman.core.embedding_indexer as embedding_indexer  # noqa: E402
import wingman.core.embedding_storage as embedding_storage  # noqa: E402
import wingman.shared.source_registry as source_registry  # noqa: E402
from products.atlas.batch_ingestion import (  # noqa: E402
    execute_batch,
    folder_file_inputs,
    preview_batch,
    resume_plan,
    write_manifest,
)
from products.atlas.intake_service import create_source_id  # noqa: E402
from products.atlas.product_config import create_atlas_context  # noqa: E402
from wingman.core.folder_intake import collect_folder_entries  # noqa: E402
from wingman.core.knowledge import retrieve_evidence  # noqa: E402


SUPPORTED_CYCLE = (".pdf", ".docx", ".xlsx", ".pptx", ".csv", ".txt", ".md", ".markdown")


def write_fixture(path, index):
    """Write one tiny, deterministic source in the extension's real format."""
    path.parent.mkdir(parents=True, exist_ok=True)
    sentinel = f"bulk sentinel {index:04d}"
    extension = path.suffix.lower()
    if extension == ".pdf":
        import pymupdf

        document = pymupdf.open()
        page = document.new_page()
        page.insert_text((72, 72), sentinel)
        document.save(path)
        document.close()
    elif extension == ".docx":
        from docx import Document

        document = Document()
        document.add_heading(f"Fixture {index}", level=1)
        document.add_paragraph(sentinel)
        document.save(path)
    elif extension == ".xlsx":
        from openpyxl import Workbook

        workbook = Workbook()
        worksheet = workbook.active
        worksheet.title = "Evidence"
        worksheet.append(["Fixture", "Evidence"])
        worksheet.append([index, sentinel])
        workbook.save(path)
        workbook.close()
    elif extension == ".pptx":
        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = f"Fixture {index}"
        slide.placeholders[1].text = sentinel
        presentation.save(path)
    elif extension == ".csv":
        path.write_text(
            f"Fixture,Evidence\n{index},{sentinel}\n",
            encoding="utf-8",
        )
    elif extension == ".txt":
        path.write_text(f"Fixture {index}\n{sentinel}\n", encoding="utf-8")
    elif extension in {".md", ".markdown"}:
        path.write_text(
            f"# Fixture {index}\n\n{sentinel}\n",
            encoding="utf-8",
        )
    else:
        raise ValueError(f"Unsupported generated fixture extension: {extension}")
    return path


def generate_corpus(root, document_count, *, include_edge_cases):
    """Generate exactly the requested number of local synthetic files."""
    if document_count < (16 if include_edge_cases else 8):
        raise ValueError("The requested corpus is too small for mixed-format coverage.")
    inputs = root / "inputs"
    inputs.mkdir(parents=True, exist_ok=True)
    special_count = 7 if include_edge_cases else 0
    regular_count = document_count - special_count
    regular_paths = []
    for index in range(regular_count):
        extension = SUPPORTED_CYCLE[index % len(SUPPORTED_CYCLE)]
        path = inputs / f"regular-{index:04d}{extension}"
        regular_paths.append(write_fixture(path, index))

    details = {}
    if include_edge_cases:
        duplicate = inputs / f"duplicate-copy{regular_paths[0].suffix}"
        duplicate.write_bytes(regular_paths[0].read_bytes())
        write_fixture(inputs / "revision-a" / "same-name.txt", 9001)
        write_fixture(inputs / "revision-b" / "same-name.txt", 9002)
        (inputs / "unsupported.rtf").write_text("unsupported", encoding="utf-8")
        (inputs / "empty.txt").write_bytes(b"")

        import pymupdf

        no_text = inputs / "no-text.pdf"
        document = pymupdf.open()
        document.new_page()
        document.save(no_text)
        document.close()
        failure = inputs / "recoverable-failure.txt"
        failure.write_text("recoverable failure sentinel", encoding="utf-8")
        details = {
            "duplicate": duplicate.relative_to(inputs).as_posix(),
            "no_text": no_text.relative_to(inputs).as_posix(),
            "failure": failure.relative_to(inputs).as_posix(),
        }
    return inputs, details


def deterministic_embedding(text):
    digest = hashlib.sha256(text.encode("utf-8")).digest()
    return [byte / 255 for byte in digest[:8]]


def baseline_knowledge(path):
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(
        json.dumps(
            [
                {
                    "id": "baseline_001",
                    "document": "baseline-source",
                    "domain": "General",
                    "heading": "Baseline",
                    "section": "General",
                    "concepts": [],
                    "records": [],
                    "location": "Line 1",
                    "text": "baseline retrieval sentinel",
                }
            ]
        ),
        encoding="utf-8",
    )


def timed_retrieval(plan):
    started = time.perf_counter()
    evidence = retrieve_evidence(plan)
    return {
        "elapsed_seconds": time.perf_counter() - started,
        "source_ids": sorted({item.get("source") for item in evidence}),
        "evidence_count": len(evidence),
    }


def run_offline_batch(document_count=500, *, include_edge_cases=False):
    """Run one disposable batch with no model or external service calls."""
    started = time.perf_counter()
    with tempfile.TemporaryDirectory(prefix="wingman-bulk-soak-") as directory:
        root = Path(directory)
        inputs_root, details = generate_corpus(
            root,
            document_count,
            include_edge_cases=include_edge_cases,
        )
        state = root / "state"
        uploads = root / "data" / "documents" / "uploads"
        manifest_path = state / "batch.json"
        baseline_knowledge(root / "data" / "documents" / "baseline.json")
        before = None
        after = None
        first_counts = None
        failed_artifacts_absent = None
        no_text_artifacts_absent = None

        environment = {
            "PYTHON_DOTENV_DISABLED": "1",
            "OPENAI_API_KEY": "bulk-ingestion-offline-placeholder",
            "WINGMAN_LEDGER_PATH": str(state / "ledger.sqlite3"),
        }
        with (
            patch.dict(os.environ, environment, clear=False),
            patch.object(intake_service, "UPLOADS_DIRECTORY", uploads),
            patch.object(
                embedding_storage,
                "EMBEDDINGS_PATH",
                state / "embeddings.json",
            ),
            patch.object(
                concept_registry_storage,
                "REGISTRY_PATH",
                state / "concepts.json",
            ),
            patch.object(
                source_registry,
                "SOURCE_REGISTRY_PATH",
                state / "legacy-source-registry.json",
            ),
            patch.object(
                concept_enrichment,
                "enrich_concepts",
                side_effect=lambda item: item,
            ),
            patch.object(
                embedding_indexer,
                "create_embedding",
                side_effect=deterministic_embedding,
            ),
            patch.object(
                intake_service.source_summary_service,
                "_response_client",
                side_effect=(
                    intake_service.source_summary_service.SummaryGenerationError(
                        "Offline soak does not call a model service."
                    )
                ),
            ),
            chdir(root),
        ):
            before = timed_retrieval(
                {
                    "text_search_terms": ["baseline retrieval sentinel"],
                    "record_search_terms": [],
                    "record_types": [],
                    "record_filters": [],
                }
            )
            entries = collect_folder_entries(inputs_root, recursive=True)
            runtime_inputs = folder_file_inputs(entries)
            override_path = next(
                item.relative_path
                for item in runtime_inputs
                if item.relative_path.startswith("regular-0001")
            )
            overrides = {override_path: "AI-202"}
            plan = preview_batch(
                runtime_inputs,
                product_context=create_atlas_context(),
                input_mode="folder",
                default_course_id="AI-101",
                course_overrides=overrides,
                product_metadata={
                    "program": "MSAIB",
                    "academic_year": "2026-2027",
                },
                assignments_confirmed=True,
                batch_id=f"batch-offline-{document_count}",
            )

            resumable = next(
                record
                for record in plan.manifest["files"]
                if record["terminal_result"] is None
            )
            resumable["attempt_count"] = 1
            resumable["progress_stage"] = "indexing"
            resumable_path = resumable["relative_path"]
            write_manifest(plan.manifest, manifest_path)
            plan = resume_plan(
                manifest_path,
                runtime_inputs,
                product_context=create_atlas_context(),
            )

            failed_once = False

            def controlled_ingestor(**kwargs):
                nonlocal failed_once
                if (
                    include_edge_cases
                    and kwargs["file_name"] == "recoverable-failure.txt"
                    and not failed_once
                ):
                    failed_once = True
                    error = RuntimeError("injected recoverable failure")
                    error.cleanup_verified = True
                    error.failure_stage = "extracting"
                    raise error
                return intake_service.ingest_uploaded_document(**kwargs)

            first_result = execute_batch(
                plan,
                product_context=create_atlas_context(),
                manifest_path=manifest_path,
                ingestor=controlled_ingestor,
            )
            first_counts = dict(first_result["counts"])

            if include_edge_cases:
                failure_bytes = (inputs_root / details["failure"]).read_bytes()
                failure_id = create_source_id(
                    "recoverable-failure.txt",
                    hashlib.sha256(failure_bytes).hexdigest(),
                )
                no_text_bytes = (inputs_root / details["no_text"]).read_bytes()
                no_text_id = create_source_id(
                    "no-text.pdf",
                    hashlib.sha256(no_text_bytes).hexdigest(),
                )
                failed_artifacts_absent = not (uploads / failure_id).exists()
                no_text_artifacts_absent = not (uploads / no_text_id).exists()
                final_result = execute_batch(
                    plan,
                    product_context=create_atlas_context(),
                    manifest_path=manifest_path,
                    retry_failed=True,
                    ingestor=intake_service.ingest_uploaded_document,
                )
            else:
                final_result = first_result

            after = timed_retrieval(
                {
                    "text_search_terms": [
                        "baseline retrieval sentinel",
                        "bulk sentinel 0001",
                    ],
                    "record_search_terms": [],
                    "record_types": [],
                    "record_filters": [],
                }
            )
            registry = source_registry.load_source_registry()
            course_counts = {}
            for metadata in registry.values():
                course_id = metadata.get("course_id")
                if course_id:
                    course_counts[course_id] = course_counts.get(course_id, 0) + 1
            revision_records = [
                record
                for record in final_result["manifest"]["files"]
                if record.get("possible_revision_of")
            ]
            retry_records = [
                record
                for record in final_result["manifest"]["files"]
                if record["visible_name"] == "recoverable-failure.txt"
            ]
            resumed_record = next(
                record
                for record in final_result["manifest"]["files"]
                if record["relative_path"] == resumable_path
            )
            summary = {
                "document_count": document_count,
                "first_counts": first_counts,
                "final_counts": dict(final_result["counts"]),
                "elapsed_seconds": time.perf_counter() - started,
                "failures": [
                    record["relative_path"]
                    for record in final_result["manifest"]["files"]
                    if record["terminal_result"] == "failed"
                ],
                "course_counts": course_counts,
                "possible_revision_count": len(revision_records),
                "resumed_attempt_count": resumed_record["attempt_count"],
                "retry_attempt_count": (
                    retry_records[0]["attempt_count"] if retry_records else None
                ),
                "failed_artifacts_absent": failed_artifacts_absent,
                "no_text_artifacts_absent": no_text_artifacts_absent,
                "retrieval_before": before,
                "retrieval_after": after,
                "manifest_version": final_result["manifest"]["manifest_version"],
                "cleanup_failure_stopped_batch": final_result["manifest"][
                    "cleanup_failure_stopped_batch"
                ],
            }
    return summary


def main(arguments=None):
    parser = argparse.ArgumentParser(
        description="Run the disposable offline mixed-format batch soak."
    )
    parser.add_argument("--count", type=int, default=500)
    parser.add_argument("--edge-cases", action="store_true")
    options = parser.parse_args(arguments)
    summary = run_offline_batch(
        options.count,
        include_edge_cases=options.edge_cases,
    )
    print(json.dumps(summary, indent=2, sort_keys=True))
    return 1 if summary["failures"] else 0


if __name__ == "__main__":
    raise SystemExit(main())
