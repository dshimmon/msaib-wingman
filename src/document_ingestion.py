import json
from pathlib import Path
from section_resolver import resolve_section
from concept_enrichment import enrich_concepts
from embedding_indexer import index_knowledge_objects
from pptx import Presentation

current_section = "General"

def detect_slide_heading(slide):
    candidates = []

    for shape in slide.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue

        font_sizes = []

        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size:
                        font_sizes.append(run.font.size.pt)

        largest_font = max(font_sizes) if font_sizes else 0
        text = shape.text.strip()

        candidates.append(
            {
                "text": text,
                "font_size": largest_font,
                "top": shape.top.inches,
                "word_count": len(text.split()),
            }
        )

    if not candidates:
        return None

    largest_font = max(candidate["font_size"] for candidate in candidates)

    heading_candidates = [
        candidate
        for candidate in candidates
        if candidate["font_size"] == largest_font
        and candidate["top"] <= 2.0
        and candidate["word_count"] <= 12
    ]

    if not heading_candidates:
        return None

    best_candidate = min(
        heading_candidates,
        key=lambda candidate: (
            candidate["word_count"],
            candidate["top"],
        ),
    )

    return best_candidate["text"]

def extract_powerpoint_chunks(file_path, domain):
    presentation = Presentation(file_path)
    chunks = []
    current_section = "General"

    for slide_number, slide in enumerate(presentation.slides, start=1):
        detected_heading = detect_slide_heading(slide)
        
        current_section = resolve_section(
            detected_heading,
            current_section,
        )

        slide_text = []

        for shape in slide.shapes:
            if shape.has_table:
                table_rows = []

                for row in shape.table.rows:
                    cell_values = [
                        cell.text.strip()
                        for cell in row.cells
                        if cell.text.strip()
                    ]

                    if cell_values:
                        table_rows.append(" | ".join(cell_values))

                if table_rows:
                    slide_text.append("\n".join(table_rows))

            elif hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        combined_text = "\n".join(slide_text)
        
        if combined_text:
            knowledge_object = {
                "id": f"{Path(file_path).stem}_{slide_number:03}",
                "document": Path(file_path).stem,
                "domain": domain,
                "heading": detected_heading,
                "section": current_section,
                "concepts": [],
                "records": [],
                "location": f"Slide {slide_number}",
                "text": combined_text,
            }
            knowledge_object = enrich_concepts(knowledge_object)
            
            chunks.append(knowledge_object)

    return (chunks)


def save_chunks(chunks, output_path):
    with open(output_path, "w") as file:
        json.dump(chunks, file, indent=2)


if __name__ == "__main__":
    chunks = extract_powerpoint_chunks(
        "data/documents/onboarding/msaib-onboarding-2026.pptx",
        "Onboarding"
    )

    save_chunks(
        chunks,
        "data/documents/onboarding/msaib-onboarding-2026.json"
    )

    index_knowledge_objects(chunks)

    print(f"Saved {len(chunks)} chunks.")