"""Atlas-owned syllabus detection and course-material classification."""

import csv
from dataclasses import dataclass
from itertools import islice
from pathlib import Path
import re
import tempfile
import time
import zipfile

import pymupdf
from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph
from openpyxl import load_workbook
from pptx import Presentation

from wingman.core.excel_adapter import flatten_worksheet_row
from wingman.core.word_adapter import extract_table_rows


MATERIAL_TYPE_LABELS = {
    "syllabus": "Syllabus",
    "notes": "Class notes",
    "lectures": "Class lectures",
    "homework": "Homework",
    "other": "Other",
}
MATERIAL_TYPES = tuple(MATERIAL_TYPE_LABELS)
MAX_PREVIEW_FILE_BYTES = 10 * 1024 * 1024
MAX_PREVIEW_TEXT_BYTES = 128 * 1024
MAX_PREVIEW_UNITS = 8
MAX_PREVIEW_ROWS = 80
MAX_PREVIEW_CHARACTERS = 12000
MAX_PREVIEW_SECONDS = 3.0
MAX_PREVIEW_ARCHIVE_ENTRIES = 1000
MAX_PREVIEW_UNCOMPRESSED_BYTES = 50 * 1024 * 1024

_COURSE_CODE_PATTERN = re.compile(
    r"\b([A-Za-z]{2,8})\s*[- ]?\s*(\d{3,5}[A-Za-z]?)\b"
)
_COURSE_TITLE_LABEL = re.compile(
    r"^(?:course\s+)?(?:title|name)\s*[:\-]\s*(.+)$",
    re.IGNORECASE,
)
_COURSE_LABEL = re.compile(
    r"^course\s*[:\-]\s*(.+)$",
    re.IGNORECASE,
)
_COURSE_CODE_LABEL = re.compile(
    r"^(?:course\s+)?(?:number|no\.?|code)\s*[:\-]\s*(.+)$",
    re.IGNORECASE,
)
_TERM_PATTERN = re.compile(
    r"\b(?:fall|spring|summer|winter|january|february|march|april|may|june|"
    r"july|august|september|october|november|december)\s+20\d{2}\b",
    re.IGNORECASE,
)
_TERM_WORDS = frozenset(
    {
        "fall",
        "spring",
        "summer",
        "winter",
        "january",
        "february",
        "march",
        "april",
        "may",
        "june",
        "july",
        "august",
        "september",
        "october",
        "november",
        "december",
    }
)
_GENERIC_TITLE_TERMS = frozenset(
    {
        "course syllabus",
        "syllabus",
        "course outline",
        "class syllabus",
        "instructor information",
        "course information",
        "course code",
        "course name",
        "course number",
        "course title",
    }
)


@dataclass(frozen=True)
class SyllabusAnalysis:
    """Best-effort, reviewable metadata inferred from one uploaded document."""

    is_syllabus: bool
    course_id: str | None
    course_name: str | None
    material_type: str


def normalize_material_type(value):
    """Normalize one stable Atlas course-material category."""
    if value is None:
        return "other"
    if not isinstance(value, str):
        raise ValueError("Material type must be text.")
    normalized = value.strip().lower()
    if not normalized:
        return "other"
    if normalized not in MATERIAL_TYPE_LABELS:
        raise ValueError(
            "Material type must be syllabus, notes, lectures, homework, or other."
        )
    return normalized


def material_type_for_catalog(value):
    """Read legacy or malformed catalog metadata without hiding the source."""
    try:
        return normalize_material_type(value)
    except (TypeError, ValueError):
        return "other"


def normalize_course_name(value):
    """Keep a bounded, readable Atlas course-folder label."""
    if value is None:
        return None
    if not isinstance(value, str):
        raise ValueError("Course name must be text.")
    if any(ord(character) < 32 or ord(character) == 127 for character in value):
        raise ValueError("Course name cannot contain control characters.")
    normalized = re.sub(r"\s+", " ", value).strip(" \t\r\n:|-/")
    if not normalized:
        return None
    if len(normalized) > 160:
        raise ValueError("Course name must be 160 characters or fewer.")
    return normalized


def _early_text(
    units,
    *,
    character_limit=MAX_PREVIEW_CHARACTERS,
    unit_limit=MAX_PREVIEW_UNITS,
    deadline=None,
):
    parts = []
    remaining = character_limit
    for unit in islice(units or (), unit_limit):
        if deadline is not None and time.monotonic() >= deadline:
            break
        heading = str(unit.get("heading") or "").strip()
        text = str(unit.get("text") or "").strip()
        for value in (heading, text):
            if not value or remaining <= 0:
                continue
            fragment = value[:remaining]
            parts.append(fragment)
            remaining -= len(fragment)
    return "\n".join(parts)


def _text_preview_units(file_path, *, unit_limit, deadline):
    with Path(file_path).open("rb") as source:
        raw_text = source.read(MAX_PREVIEW_TEXT_BYTES + 1)
    raw_text = raw_text[:MAX_PREVIEW_TEXT_BYTES]
    text = raw_text.decode("utf-8-sig")
    units = []
    group = []
    for line in text.splitlines():
        if time.monotonic() >= deadline or len(units) >= unit_limit:
            break
        if line.strip():
            group.append(line.rstrip())
        elif group:
            units.append({"heading": None, "text": "\n".join(group)})
            group = []
    if group and len(units) < unit_limit and time.monotonic() < deadline:
        units.append({"heading": None, "text": "\n".join(group)})
    return units


def _csv_preview_units(file_path, *, unit_limit, deadline):
    units = []
    with Path(file_path).open(
        "r", encoding="utf-8-sig", newline=""
    ) as source:
        rows = csv.reader(source, strict=True)
        for row_number, row in enumerate(
            islice(rows, unit_limit), start=1
        ):
            if time.monotonic() >= deadline:
                break
            if any(value.strip() for value in row):
                units.append(
                    {
                        "heading": "CSV columns" if row_number == 1 else None,
                        "text": " | ".join(row),
                    }
                )
    return units


def _pdf_preview_units(file_path, *, unit_limit, deadline):
    units = []
    with pymupdf.open(file_path) as document:
        page_limit = min(document.page_count, unit_limit)
        for page_number in range(page_limit):
            if time.monotonic() >= deadline:
                break
            page_text = document.load_page(page_number).get_text(
                "text", sort=True
            ).strip()
            if page_text:
                units.append({"heading": None, "text": page_text})
    return units


def _powerpoint_preview_units(file_path, *, unit_limit, deadline):
    units = []
    presentation = Presentation(file_path)
    for slide in islice(presentation.slides, unit_limit):
        if time.monotonic() >= deadline:
            break
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
            if sum(len(part) for part in text_parts) >= MAX_PREVIEW_CHARACTERS:
                break
        if text_parts:
            units.append({"heading": None, "text": "\n".join(text_parts)})
    return units


def _word_preview_units(file_path, *, unit_limit, deadline):
    units = []
    document = Document(file_path)
    for block in islice(document.iter_inner_content(), unit_limit):
        if time.monotonic() >= deadline:
            break
        if isinstance(block, Paragraph):
            block_text = block.text.strip()
        elif isinstance(block, Table):
            block_text = "\n".join(extract_table_rows(block)).strip()
        else:
            block_text = ""
        if block_text:
            units.append({"heading": None, "text": block_text})
    return units


def _excel_preview_units(file_path, *, unit_limit, deadline):
    units = []
    rows_scanned = 0
    workbook = load_workbook(file_path, data_only=False, read_only=True)
    try:
        for worksheet in workbook.worksheets:
            for row in worksheet.iter_rows():
                if (
                    time.monotonic() >= deadline
                    or len(units) >= unit_limit
                    or rows_scanned >= MAX_PREVIEW_ROWS
                ):
                    return units
                rows_scanned += 1
                row_text = flatten_worksheet_row(row)
                if row_text:
                    units.append(
                        {"heading": worksheet.title, "text": row_text}
                    )
    finally:
        workbook.close()
    return units


def _validate_office_preview_archive(file_path):
    with zipfile.ZipFile(file_path) as archive:
        members = archive.infolist()
        if len(members) > MAX_PREVIEW_ARCHIVE_ENTRIES:
            raise ValueError("Office preview archive has too many entries.")
        if (
            sum(member.file_size for member in members)
            > MAX_PREVIEW_UNCOMPRESSED_BYTES
        ):
            raise ValueError("Office preview archive expands beyond its limit.")


def _bounded_preview_units(file_path, *, unit_limit, deadline):
    extension = Path(file_path).suffix.lower()
    extractors = {
        ".pdf": _pdf_preview_units,
        ".pptx": _powerpoint_preview_units,
        ".docx": _word_preview_units,
        ".xlsx": _excel_preview_units,
        ".csv": _csv_preview_units,
        ".txt": _text_preview_units,
        ".md": _text_preview_units,
        ".markdown": _text_preview_units,
    }
    extractor = extractors.get(extension)
    if extractor is None:
        return ()
    if extension in {".pptx", ".docx", ".xlsx"}:
        _validate_office_preview_archive(file_path)
    return extractor(
        file_path,
        unit_limit=unit_limit,
        deadline=deadline,
    )


def _syllabus_score(file_name, text):
    stem = Path(file_name).stem.replace("_", " ").replace("-", " ")
    filename_text = re.sub(r"\s+", " ", stem).lower()
    early_text = text.lower()
    score = 0
    if "syllabus" in filename_text or "course outline" in filename_text:
        score += 4
    if re.search(
        r"(?mi)^(?:course\s+|class\s+)?(?:syllabus|course\s+outline)\b",
        text,
    ):
        score += 3
    elif re.search(r"\bsyllabus\b", early_text):
        score += 1
    for phrase in (
        "course description",
        "learning objectives",
        "instructor",
        "grading policy",
        "office hours",
        "required materials",
    ):
        if phrase in early_text:
            score += 1
    return score


def classify_material_type(file_name, text=""):
    """Classify a material from conservative filename and early-text signals."""
    normalized_name = re.sub(
        r"[^a-z0-9]+", " ", Path(file_name).stem.lower()
    ).strip()
    normalized_text = text.lower()
    if _syllabus_score(file_name, text) >= 3:
        return "syllabus"
    if re.search(r"\b(?:homework|assignment|problem\s+set|hw\s*\d*)\b", normalized_name):
        return "homework"
    if re.search(r"\b(?:lecture|slides?|deck|presentation)\b", normalized_name):
        return "lectures"
    if re.search(r"\b(?:class\s+notes?|study\s+notes?|notes?)\b", normalized_name):
        return "notes"
    if "homework" in normalized_text[:2000] and "due" in normalized_text[:2000]:
        return "homework"
    return "other"


def _normalize_course_code(value):
    for match in _COURSE_CODE_PATTERN.finditer(value or ""):
        subject = match.group(1).upper()
        number = match.group(2).upper()
        if (
            subject.casefold() in _TERM_WORDS
            and re.fullmatch(r"(?:19|20)\d{2}", number)
        ):
            continue
        return f"{subject} {number}"
    return None


def _clean_title(value):
    if not value:
        return None
    title = _TERM_PATTERN.sub("", value)
    title = re.sub(r"\b(?:course\s+)?syllabus\b", "", title, flags=re.IGNORECASE)
    title = re.sub(r"\bcourse\s+outline\b", "", title, flags=re.IGNORECASE)
    title = _COURSE_CODE_PATTERN.sub("", title, count=1)
    title = re.sub(r"\s+", " ", title).strip(" \t\r\n:|-/–—()[]")
    if not title or title.casefold() in _GENERIC_TITLE_TERMS:
        return None
    if len(title) > 160 or len(title.split()) > 20:
        return None
    if any(
        marker in title.casefold()
        for marker in ("@", "instructor", "office hours", "department of")
    ):
        return None
    return normalize_course_name(title)


def _safe_course_id_from_name(course_name):
    safe = re.sub(r"[^A-Za-z0-9._:/ -]+", " ", course_name or "")
    safe = re.sub(r"\s+", " ", safe).strip(" ._:/-")
    return safe[:120].rstrip() or None


def _course_identity(text, file_name):
    lines = [re.sub(r"\s+", " ", line).strip() for line in text.splitlines()]
    lines = [line for line in lines if line][:100]
    course_code = None
    course_name = None

    for line in lines:
        code_label = _COURSE_CODE_LABEL.match(line)
        if code_label and course_code is None:
            course_code = _normalize_course_code(code_label.group(1))
        title_label = _COURSE_TITLE_LABEL.match(line)
        if title_label and course_name is None:
            course_name = _clean_title(title_label.group(1))
        course_label = _COURSE_LABEL.match(line)
        if course_label:
            value = course_label.group(1)
            course_code = course_code or _normalize_course_code(value)
            course_name = course_name or _clean_title(value)

    for index, line in enumerate(lines):
        opening_match = _COURSE_CODE_PATTERN.match(line)
        candidate_code = _normalize_course_code(
            opening_match.group(0) if opening_match else None
        )
        if candidate_code is None:
            continue
        remainder = _clean_title(line)
        if remainder:
            course_code = course_code or candidate_code
            course_name = course_name or remainder
            break
        course_code = course_code or candidate_code
        if course_name is None:
            for neighbor in lines[index + 1 : index + 3]:
                candidate_title = _clean_title(neighbor)
                if candidate_title and _normalize_course_code(neighbor) is None:
                    course_name = candidate_title
                    break

    filename_text = re.sub(r"[_-]+", " ", Path(file_name).stem)
    filename_code = _normalize_course_code(filename_text)
    filename_name = _clean_title(filename_text)
    course_code = course_code or filename_code
    if (
        course_name is None
        and filename_name
        and (filename_code or len(filename_name.split()) >= 2)
    ):
        course_name = filename_name

    if course_name is None:
        for index, line in enumerate(lines):
            if line.casefold() in _GENERIC_TITLE_TERMS:
                neighbors = lines[max(0, index - 2) : index] + lines[index + 1 : index + 3]
                course_name = next(
                    (candidate for candidate in map(_clean_title, neighbors) if candidate),
                    None,
                )
                if course_name:
                    break

    course_name = course_name or filename_name
    course_id = course_code or _safe_course_id_from_name(course_name)
    course_name = course_name or course_code
    return course_id, course_name


def analyze_uploaded_document(
    file_name,
    file_bytes,
    *,
    unit_extractor=None,
):
    """Inspect a bounded opening preview and return reviewable upload metadata."""
    safe_name = Path(file_name).name
    text = ""
    if file_bytes and len(file_bytes) <= MAX_PREVIEW_FILE_BYTES:
        try:
            with tempfile.TemporaryDirectory(prefix="atlas-syllabus-") as directory:
                temporary_path = Path(directory) / safe_name
                temporary_path.write_bytes(file_bytes)
                deadline = time.monotonic() + MAX_PREVIEW_SECONDS
                units = (
                    unit_extractor(temporary_path)
                    if unit_extractor is not None
                    else _bounded_preview_units(
                        temporary_path,
                        unit_limit=MAX_PREVIEW_UNITS,
                        deadline=deadline,
                    )
                )
                text = _early_text(units, deadline=deadline)
        except Exception:
            # Preview inference is advisory; normal ingestion reports document errors.
            text = ""

    material_type = classify_material_type(safe_name, text)
    is_syllabus = material_type == "syllabus"
    course_id, course_name = (
        _course_identity(text, safe_name) if is_syllabus else (None, None)
    )
    return SyllabusAnalysis(
        is_syllabus=is_syllabus,
        course_id=course_id,
        course_name=course_name,
        material_type=material_type,
    )
