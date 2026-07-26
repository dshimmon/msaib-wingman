# Extracts normalized document units from PowerPoint files.

from pptx import Presentation


def detect_slide_heading(slide):
    """
    Identify the most likely heading on a slide.
    """
    candidates = []

    for shape in slide.shapes:
        if not hasattr(shape, "text") or not shape.text.strip():
            continue

        font_sizes = []

        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size:
                        font_sizes.append(run.font.size.pt)

        candidates.append(
            {
                "text": shape.text.strip(),
                "font_size": (
                    max(font_sizes)
                    if font_sizes
                    else 0
                ),
                "top": shape.top.inches,
                "word_count": len(
                    shape.text.strip().split()
                ),
            }
        )

    if not candidates:
        return None

    largest_font = max(
        candidate["font_size"]
        for candidate in candidates
    )

    heading_candidates = [
        candidate
        for candidate in candidates
        if candidate["font_size"] == largest_font
        and candidate["top"] <= 2.0
        and candidate["word_count"] <= 12
    ]

    if not heading_candidates:
        return None

    best_candidate = min(
        heading_candidates,
        key=lambda candidate: (
            candidate["word_count"],
            candidate["top"],
        ),
    )

    return best_candidate["text"]


def extract_powerpoint_units(file_path):
    """
    Extract one normalized document unit per slide.
    """
    presentation = Presentation(file_path)
    units = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        slide_text = []

        for shape in slide.shapes:
            if shape.has_table:
                table_rows = []

                for row in shape.table.rows:
                    cell_values = [
                        cell.text.strip()
                        for cell in row.cells
                        if cell.text.strip()
                    ]

                    if cell_values:
                        table_rows.append(
                            " | ".join(cell_values)
                        )

                if table_rows:
                    slide_text.append(
                        "\n".join(table_rows)
                    )

            elif (
                hasattr(shape, "text")
                and shape.text.strip()
            ):
                slide_text.append(
                    shape.text.strip()
                )

        combined_text = "\n".join(slide_text)

        if combined_text:
            units.append(
                {
                    "heading": detect_slide_heading(
                        slide
                    ),
                    "location": (
                        f"Slide {slide_number}"
                    ),
                    "text": combined_text,
                }
            )

    return units